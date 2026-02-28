import streamlit as st
import google.generativeai as genai
import tempfile
from docx import Document
from pptx import Presentation

# 1. 頁面設定
st.set_page_config(page_title="AI 智能摘要工具", layout="wide")
st.title("🌟 多媒體文件摘要與翻譯器")

# 2. 直接顯示上傳介面 (不需要 Key 就能看到)
uploaded_file = st.file_uploader("📂 步驟 1：請拖曳或選擇要摘要的檔案 (PPT, Word, MP3, MP4)", type=["docx", "pptx", "mp3", "mp4"])

# 3. 側邊欄：輸入 API Key
with st.sidebar:
    st.header("🔑 步驟 2：金鑰設定")
    api_key = st.text_input("請輸入 Google API Key 以啟動 AI", type="password")
    st.info("支援格式：\n- Word (.docx)\n- PPT (.pptx)\n- 影音 (.mp3, .mp4)")

# 4. 核心處理邏輯
if uploaded_file:
    if not api_key:
        st.warning("⚠️ 偵測到已上傳檔案，請在左側選單輸入 **Google API Key** 以開始 AI 分析。")
    else:
        with st.spinner("🚀 AI 正在深度分析中，請稍候..."):
            try:
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel('gemini-1.5-flash')
                
                # 處理影音
                if uploaded_file.type in ["audio/mpeg", "video/mp4"]:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}") as f:
                        f.write(uploaded_file.getbuffer())
                        g_file = genai.upload_file(path=f.name)
                        response = model.generate_content([g_file, "請擷取此內容重點摘要，並用繁體中文列出"])
                
                # 處理文件
                else:
                    text_content = ""
                    if "word" in uploaded_file.type:
                        doc = Document(uploaded_file)
                        text_content = "\n".join([p.text for p in doc.paragraphs])
                    elif "presentation" in uploaded_file.type:
                        prs = Presentation(uploaded_file)
                        text_content = "\n".join([s.text for slide in prs.slides for s in slide.shapes if hasattr(s, "text")])
                    
                    response = model.generate_content(f"請將以下內容整理成重點摘要，並翻譯為繁體中文：\n\n{text_content}")

                st.success("✨ 分析完成！")
                st.subheader("📝 重點摘要結果")
                st.markdown(response.text)
                
            except Exception as e:
                st.error(f"❌ 發生錯誤：{e}")
